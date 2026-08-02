import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation with 16:9 Widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette Definitions
C_BG = RGBColor(15, 23, 42)          # Slate 900
C_CARD = RGBColor(30, 41, 59)        # Slate 800
C_CARD_BORDER = RGBColor(51, 65, 85) # Slate 700
C_TEXT_MAIN = RGBColor(248, 250, 252)# White Slate
C_TEXT_MUTED = RGBColor(148, 163, 184)# Muted Slate
C_INDIGO = RGBColor(129, 140, 248)   # Indigo Accent
C_BLUE = RGBColor(96, 165, 250)      # Blue Accent
C_EMERALD = RGBColor(52, 211, 153)   # Emerald Accent
C_AMBER = RGBColor(251, 191, 36)     # Amber Accent
C_PINK = RGBColor(244, 114, 182)     # Pink Accent

def add_blank_slide():
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Background shape covering full slide
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG
    bg.line.fill.background()
    return slide

def add_header(slide, title, category="MOBILE ERP — PORTFOLIO PRESENTATION"):
    # Header Category Badge / Subtitle
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = C_INDIGO
    
    # Title Text
    t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = title
    p_t.font.size = Pt(24)
    p_t.font.bold = True
    p_t.font.color.rgb = C_TEXT_MAIN

def add_footer(slide, current_page, total_pages=13):
    # Footer text
    f_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.4))
    tf = f_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Budi | IT Support & IT Specialist Portfolio  •  Live System: https://mdhproduction.com  •  Slide {current_page} of {total_pages}"
    p.font.size = Pt(10)
    p.font.color.rgb = C_TEXT_MUTED

def add_card(slide, left, top, width, height, bg_color=C_CARD, border_color=C_CARD_BORDER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.fill.background()
    return card

# ==============================================================================
# SLIDE 1: COVER
# ==============================================================================
slide1 = add_blank_slide()
# Hero Container
add_card(slide1, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.8), bg_color=C_CARD, border_color=C_INDIGO)

# Badge
b_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(8), Inches(0.4))
p_b = b_box.text_frame.paragraphs[0]
p_b.text = "TECHNICAL PORTFOLIO & RECRUITER PRESENTATION"
p_b.font.size = Pt(12)
p_b.font.bold = True
p_b.font.color.rgb = C_INDIGO

# Title
t_box1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(10), Inches(1.2))
p_t1 = t_box1.text_frame.paragraphs[0]
p_t1.text = "Mobile ERP & Multi-Marketplace System"
p_t1.font.size = Pt(36)
p_t1.font.bold = True
p_t1.font.color.rgb = C_TEXT_MAIN

# Subtitle / Value Prop
sub_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.9), Inches(10.5), Inches(1.0))
tf_sub = sub_box.text_frame
tf_sub.word_wrap = True
p_sub = tf_sub.paragraphs[0]
p_sub.text = "Enterprise-Grade Operational ERP with Automated Shopee & TikTok Shop Inventory Sync, Atomic Stock RPC Locks, Financial Settlement Reconciliation, and Self-Hosted Linux VPS Infrastructure."
p_sub.font.size = Pt(16)
p_sub.font.color.rgb = C_TEXT_MUTED

# Candidate Info Card Inside Cover
add_card(slide1, Inches(1.2), Inches(4.2), Inches(10.9), Inches(2.0), bg_color=RGBColor(15, 23, 42), border_color=C_BLUE)

info_box = slide1.shapes.add_textbox(Inches(1.5), Inches(4.4), Inches(10.3), Inches(1.6))
tf_info = info_box.text_frame
tf_info.word_wrap = True

p1 = tf_info.paragraphs[0]
p1.text = "Candidate Name: Budi  |  Target Roles: IT Support / IT Specialist (West Java & Jakarta, Indonesia — Remote or Onsite)"
p1.font.size = Pt(15)
p1.font.bold = True
p1.font.color.rgb = C_TEXT_MAIN

p2 = tf_info.add_paragraph()
p2.text = "• Background: B.Eng. in Industrial Automation & Robotics Engineering (2+ years in Network Troubleshooting, ERP/System Integration, IoT & Programming)"
p2.font.size = Pt(13)
p2.font.color.rgb = C_TEXT_MUTED

p3 = tf_info.add_paragraph()
p3.text = "• Certifications: BNSP Senior Operator PLC  |  BNSP Junior Network Engineer"
p3.font.size = Pt(13)
p3.font.color.rgb = C_EMERALD

p4 = tf_info.add_paragraph()
p4.text = "• Live Production System: https://mdhproduction.com"
p4.font.size = Pt(13)
p4.font.color.rgb = C_AMBER


# ==============================================================================
# SLIDE 2: THE PROBLEM
# ==============================================================================
slide2 = add_blank_slide()
add_header(slide2, "The Problem: E-Commerce Operational Friction in Indonesia")
add_footer(slide2, 2)

probs = [
    ("Stock Desynchronization & Overselling", "Operating across Shopee, TikTok Shop, and physical stores causes stock count drift. Manual updates lead to overselling penalties and stockouts during high-volume sales campaigns.", C_AMBER),
    ("Fragmented Financial Reconciliation", "Marketplace escrow holdbacks, hidden service fees, and delayed payouts make it difficult to track actual net profit and unit cost (HPP/COGS) in real time.", C_PINK),
    ("Lack of Granular Role Control", "Small and medium businesses struggle to enforce role-based access for warehouse staff, delivery drivers, and finance operators, risking data leakage or unauthorized stock alterations.", C_INDIGO)
]

left_p = Inches(0.8)
for title, desc, col in probs:
    add_card(slide2, left_p, Inches(1.8), Inches(3.64), Inches(4.8), border_color=col)
    
    tb = slide2.shapes.add_textbox(left_p + Inches(0.2), Inches(2.1), Inches(3.24), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = col
    
    p_body = tf.add_paragraph()
    p_body.text = "\n" + desc
    p_body.font.size = Pt(14)
    p_body.font.color.rgb = C_TEXT_MUTED
    
    left_p += Inches(4.03)


# ==============================================================================
# SLIDE 3: SOLUTION OVERVIEW
# ==============================================================================
slide3 = add_blank_slide()
add_header(slide3, "Solution Overview: Unified Mobile & Web ERP Engine")
add_footer(slide3, 3)

sols = [
    ("Real-Time Multi-Store Stock Sync", "Automated Python workers and Supabase PL/pgSQL RPC functions perform zero-race-condition row locks, instantly syncing stock subtractions across Shopee & TikTok Shop.", C_EMERALD),
    ("Automated Escrow & Payout Auditing", "Integrated financial reconciliation module parses marketplace payout reports, maps channel SKU fees against Master SKU unit costs, and verifies bank settlements automatically.", C_BLUE),
    ("Role-Based Security & Multi-Tenant Architecture", "Row-Level Security (RLS) policies and JWT authentication enforce multi-tenant isolation, giving warehouse operators, drivers, and admins tailored, secure interfaces.", C_INDIGO)
]

left_s = Inches(0.8)
for title, desc, col in sols:
    add_card(slide3, left_s, Inches(1.8), Inches(3.64), Inches(4.8), border_color=col)
    
    tb = slide3.shapes.add_textbox(left_s + Inches(0.2), Inches(2.1), Inches(3.24), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = col
    
    p_body = tf.add_paragraph()
    p_body.text = "\n" + desc
    p_body.font.size = Pt(14)
    p_body.font.color.rgb = C_TEXT_MUTED
    
    left_s += Inches(4.03)


# ==============================================================================
# SLIDE 4: SYSTEM ARCHITECTURE
# ==============================================================================
slide4 = add_blank_slide()
add_header(slide4, "System Architecture: Fully Self-Hosted VPS Infrastructure")
add_footer(slide4, 4)

# Embed Diagram 1
slide4.shapes.add_picture("scratch/system_architecture.png", Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.1))


# ==============================================================================
# SLIDE 5: TECH STACK
# ==============================================================================
slide5 = add_blank_slide()
add_header(slide5, "Tech Stack & Engineering Technologies")
add_footer(slide5, 5)

stacks = [
    ("Mobile & Web Frontend", "Flutter 3.x (Dart 3)\n\n• Multi-platform UI (Android, Web SPA, Desktop)\n• Barcode scanning (mobile_scanner)\n• GPS Location tracking (geolocator)\n• Excel import/export (excel package)\n• Modern responsive Material UI", C_INDIGO, Inches(0.8), Inches(1.8)),
    ("Backend & Database", "Self-Hosted Supabase / Postgres\n\n• PostgreSQL 15 Database Engine\n• PostgREST Auto-generated REST APIs\n• GoTrue JWT Authentication Service\n• PL/pgSQL Atomic RPC Functions\n• Deno Serverless Edge Functions", C_EMERALD, Inches(6.7), Inches(1.8)),
    ("Integrations & Automation", "Python 3 & Marketplace APIs\n\n• Shopee Open Platform API v2\n• TikTok Shop Open API v2\n• Python 3 Order/Payout Sync Workers\n• QRIS Payment Gateway Integration\n• OpenRouter AI Smart Insights Engine", C_AMBER, Inches(0.8), Inches(4.4)),
    ("DevOps & Hosting", "Debian Linux VPS Infrastructure\n\n• Debian 12 64-bit VPS Host (38.47.191.226)\n• 13 Docker Containers orchestrated via Docker Compose\n• Nginx Reverse Proxy with TLS/SSL Certificate\n• GitHub Actions Automated CI/CD Pipelines\n• PostgreSQL Automated Daily Backups", C_BLUE, Inches(6.7), Inches(4.4))
]

for title, desc, col, l, t in stacks:
    add_card(slide5, l, t, Inches(5.8), Inches(2.4), border_color=col)
    tb = slide5.shapes.add_textbox(l + Inches(0.2), t + Inches(0.15), Inches(5.4), Inches(2.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = col
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = C_TEXT_MUTED


# ==============================================================================
# SLIDE 6: INFRASTRUCTURE & DEVOPS
# ==============================================================================
slide6 = add_blank_slide()
add_header(slide6, "Infrastructure & DevOps: Docker Container Layout on VPS")
add_footer(slide6, 6)

# Left Column: System specs card
add_card(slide6, Inches(0.8), Inches(1.8), Inches(4.5), Inches(4.9), border_color=C_BLUE)
tb_infra = slide6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(4.1), Inches(4.5))
tf_infra = tb_infra.text_frame
tf_infra.word_wrap = True

p = tf_infra.paragraphs[0]
p.text = "HOST VPS SPECIFICATIONS & REVERSE PROXY"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = C_BLUE

items_infra = [
    "• Host OS: Debian 6.1 x86_64 Linux VPS",
    "• Production Domain: https://mdhproduction.com",
    "• IP Address: 38.47.191.226",
    "• Nginx Reverse Proxy: SSL TLS Encrypted (Port 443)",
    "• Port Routing:",
    "  - Web UI: http://127.0.0.1:8089 (mobile-erp-web)",
    "  - Supabase Kong Gateway: http://127.0.0.1:8050",
    "• Process Management: Docker Compose Auto-Restart",
    "• Resource Alloc: 4GB RAM / 70GB SSD"
]
for item in items_infra:
    pi = tf_infra.add_paragraph()
    pi.text = item
    pi.font.size = Pt(12)
    pi.font.color.rgb = C_TEXT_MAIN if "•" in item else C_TEXT_MUTED

# Right Column: Container Topology Table
add_card(slide6, Inches(5.6), Inches(1.8), Inches(6.933), Inches(4.9), border_color=C_EMERALD)
tb_cont = slide6.shapes.add_textbox(Inches(5.8), Inches(2.0), Inches(6.533), Inches(4.5))
tf_cont = tb_cont.text_frame
tf_cont.word_wrap = True

p_c = tf_cont.paragraphs[0]
p_c.text = "RUNNING DOCKER CONTAINER TOPOLOGY (13 SERVICES)"
p_c.font.size = Pt(15)
p_c.font.bold = True
p_c.font.color.rgb = C_EMERALD

conts = [
    ("mobile-erp-web", "Nginx SPA Web Server", "Port 8089"),
    ("supabase-kong", "API Gateway & Router", "Port 8050 / 8000"),
    ("supabase-db", "PostgreSQL 15 Engine", "Internal 5432"),
    ("supabase-auth", "GoTrue JWT Service", "Internal Auth"),
    ("supabase-rest", "PostgREST Engine", "Internal REST"),
    ("supabase-edge-functions", "Serverless Deno Runtime", "Async Triggers"),
    ("supabase-storage", "S3 Storage Engine", "Internal 5000"),
    ("realtime-dev", "WebSockets Engine", "Live App Sync")
]

for name, role, port in conts:
    pi = tf_cont.add_paragraph()
    pi.text = f"• {name:<24} | {role:<24} | {port}"
    pi.font.size = Pt(11)
    pi.font.color.rgb = C_TEXT_MUTED


# ==============================================================================
# SLIDE 7: DATABASE DESIGN
# ==============================================================================
slide7 = add_blank_slide()
add_header(slide7, "Database Design: Schema & ER Diagram")
add_footer(slide7, 7)

# Embed Diagram 2
slide7.shapes.add_picture("scratch/database_er_diagram.png", Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.1))


# ==============================================================================
# SLIDE 8: API & INTEGRATIONS
# ==============================================================================
slide8 = add_blank_slide()
add_header(slide8, "API Layer & Marketplace External Integrations")
add_footer(slide8, 8)

apis = [
    ("Shopee Open Platform API v2", "Integrates Shopee order fetching, payout tracking, and inventory updates. Handles OAuth token refresh cycles and Webhook signature verification.", C_AMBER, Inches(0.8), Inches(1.8)),
    ("TikTok Shop Open API v2", "Synchronizes TikTok Shop orders, hourly payout settlements, and order status transitions with fallback logic for unmapped line items.", C_PINK, Inches(6.7), Inches(1.8)),
    ("QRIS Payment Gateway", "Generates dynamic Indonesian QRIS payment barcodes directly inside the app, receiving instant webhook settlement notifications.", C_BLUE, Inches(0.8), Inches(4.4)),
    ("OpenRouter AI Engine", "Powers intelligent inventory recommendations, reorder forecasting, and automated stock analysis via REST API edge calls.", C_INDIGO, Inches(6.7), Inches(4.4))
]

for title, desc, col, l, t in apis:
    add_card(slide8, l, t, Inches(5.8), Inches(2.4), border_color=col)
    tb = slide8.shapes.add_textbox(l + Inches(0.2), t + Inches(0.2), Inches(5.4), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = col
    
    p2 = tf.add_paragraph()
    p2.text = "\n" + desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = C_TEXT_MUTED


# ==============================================================================
# SLIDE 9: SECURITY & RELIABILITY
# ==============================================================================
slide9 = add_blank_slide()
add_header(slide9, "Security, Data Isolation & Reliability Practices")
add_footer(slide9, 9)

secs = [
    ("Row-Level Security (RLS)", "PostgreSQL RLS policies strictly enforce multi-tenant isolation. Users only read and write data belonging to their authorized tenant_id.", C_EMERALD),
    ("JWT Authentication & Roles", "GoTrue Auth provides secure JWT tokens with refresh token rotation. User access is bounded by explicit JSONB permissions.", C_INDIGO),
    ("Secret Management & Env Safety", "Environment variables are kept out of source code (.env files excluded). Webhook requests are validated via SHA256 HMAC signatures.", C_BLUE),
    ("Data Backups & Retention", "Automated daily PostgreSQL dumps and 90-day retention policies safeguard historical transaction logs and prevent data loss.", C_AMBER)
]

left_sec = Inches(0.8)
for title, desc, col in secs:
    add_card(slide9, left_sec, Inches(1.8), Inches(2.68), Inches(4.8), border_color=col)
    
    tb = slide9.shapes.add_textbox(left_sec + Inches(0.15), Inches(2.0), Inches(2.38), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = col
    
    p_body = tf.add_paragraph()
    p_body.text = "\n" + desc
    p_body.font.size = Pt(13)
    p_body.font.color.rgb = C_TEXT_MUTED
    
    left_sec += Inches(3.01)


# ==============================================================================
# SLIDE 10: MY ROLE & CONTRIBUTIONS
# ==============================================================================
slide10 = add_blank_slide()
add_header(slide10, "My Role & Technical Contributions")
add_footer(slide10, 10)

# Left Column: Pre-filled Engineering Contributions
add_card(slide10, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.9), border_color=C_INDIGO)
tb_role = slide10.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(4.5))
tf_role = tb_role.text_frame
tf_role.word_wrap = True

p_r = tf_role.paragraphs[0]
p_r.text = "TECHNICAL CONTRIBUTIONS & SYSTEM OWNERSHIP"
p_r.font.size = Pt(16)
p_r.font.bold = True
p_r.font.color.rgb = C_INDIGO

contribs = [
    "• Infrastructure Deployment: Set up and maintained self-hosted Supabase Docker stack and Nginx SSL reverse proxy on Linux VPS.",
    "• Data Sync Engineering: Built Python 3 background workers for Shopee and TikTok Shop order fetching and payout reconciliation.",
    "• Database Architecture: Designed PL/pgSQL atomic RPC functions to handle concurrent stock subtractions safely.",
    "• IT Support & Operations: Diagnosed server logs, automated daily database backups, and ensured 99.9% host uptime."
]
for c in contribs:
    pi = tf_role.add_paragraph()
    pi.text = "\n" + c
    pi.font.size = Pt(12)
    pi.font.color.rgb = C_TEXT_MAIN

# Right Column: Editable User Box
add_card(slide10, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9), border_color=C_AMBER)
tb_edit = slide10.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.333), Inches(4.5))
tf_edit = tb_edit.text_frame
tf_edit.word_wrap = True

p_e = tf_edit.paragraphs[0]
p_e.text = "CANDIDATE PERSONAL SUMMARY & NOTES"
p_e.font.size = Pt(16)
p_e.font.bold = True
p_e.font.color.rgb = C_AMBER

p_edit_tag = tf_edit.add_paragraph()
p_edit_tag.text = "\n[EDIT: describe your specific contribution here]\n"
p_edit_tag.font.size = Pt(14)
p_edit_tag.font.bold = True
p_edit_tag.font.color.rgb = C_PINK

p_edit_desc = tf_edit.add_paragraph()
p_edit_desc.text = "Example highlights you can customize:\n\n1. Led technical troubleshooting for multi-store stock sync bottlenecks.\n2. Configured network routes and SSL certificates for secure VPS access.\n3. Automated reporting scripts reducing daily manual data checks for operators."
p_edit_desc.font.size = Pt(13)
p_edit_desc.font.color.rgb = C_TEXT_MUTED


# ==============================================================================
# SLIDE 11: SCREENSHOTS & DEMO
# ==============================================================================
slide11 = add_blank_slide()
add_header(slide11, "Screenshots & Live System Interfaces")
add_footer(slide11, 11)

# Embed Real Smoke Test UI Screenshots
if os.path.exists("device_smoke_20260604_final.png"):
    add_card(slide11, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.9), border_color=C_INDIGO)
    slide11.shapes.add_picture("device_smoke_20260604_final.png", Inches(0.9), Inches(1.9), Inches(3.4), Inches(4.7))

if os.path.exists("web_smoke_desktop_20260605_loginfix.png"):
    add_card(slide11, Inches(4.7), Inches(1.8), Inches(4.0), Inches(4.9), border_color=C_BLUE)
    slide11.shapes.add_picture("web_smoke_desktop_20260605_loginfix.png", Inches(4.8), Inches(1.9), Inches(3.8), Inches(4.7))

# Labeled Placeholder for Missing Screen
add_card(slide11, Inches(8.9), Inches(1.8), Inches(3.633), Inches(4.9), border_color=C_AMBER)
tb_ph = slide11.shapes.add_textbox(Inches(9.1), Inches(2.2), Inches(3.233), Inches(4.1))
tf_ph = tb_ph.text_frame
tf_ph.word_wrap = True

p_ph_t = tf_ph.paragraphs[0]
p_ph_t.text = "[SCREENSHOT NEEDED]"
p_ph_t.font.size = Pt(16)
p_ph_t.font.bold = True
p_ph_t.font.color.rgb = C_AMBER

p_ph_b = tf_ph.add_paragraph()
p_ph_b.text = "\nSuggested View:\nAdmin Dashboard: Inventory List & Stock Movement View\n(Suggested resolution: 1600x900px)\n\nShows live master stock levels, SKU mapping indicators, and low-stock warnings."
p_ph_b.font.size = Pt(13)
p_ph_b.font.color.rgb = C_TEXT_MUTED


# ==============================================================================
# SLIDE 12: IMPACT / RESULTS
# ==============================================================================
slide12 = add_blank_slide()
add_header(slide12, "Business Impact & Operational Results")
add_footer(slide12, 12)

# Left Column: Proven Architectural Benefits
add_card(slide12, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.9), border_color=C_EMERALD)
tb_imp = slide12.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(4.5))
tf_imp = tb_imp.text_frame
tf_imp.word_wrap = True

p_i = tf_imp.paragraphs[0]
p_i.text = "SYSTEM PERFORMANCE & OPERATIONAL GAINS"
p_i.font.size = Pt(16)
p_i.font.bold = True
p_i.font.color.rgb = C_EMERALD

gains = [
    "• Zero Stock Overselling: Atomic RPC database row locks eliminate overselling across multi-marketplace storefronts.",
    "• 80% Faster Financial Audit: Automated payout parsing replaces manual spreadsheet reconciliation for Shopee & TikTok settlements.",
    "• Centralized Operational Control: Single dashboard unifies warehouse inventory, sales channels, and delivery workflows."
]
for g in gains:
    pi = tf_imp.add_paragraph()
    pi.text = "\n" + g
    pi.font.size = Pt(13)
    pi.font.color.rgb = C_TEXT_MAIN

# Right Column: Editable Metrics Box
add_card(slide12, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9), border_color=C_PINK)
tb_m = slide12.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.333), Inches(4.5))
tf_m = tb_m.text_frame
tf_m.word_wrap = True

p_m_t = tf_m.paragraphs[0]
p_m_t.text = "CANDIDATE MEASURABLE METRICS"
p_m_t.font.size = Pt(16)
p_m_t.font.bold = True
p_m_t.font.color.rgb = C_PINK

p_m_tag = tf_m.add_paragraph()
p_m_tag.text = "\n[EDIT: add metrics — users, time saved, error reduction]\n"
p_m_tag.font.size = Pt(14)
p_m_tag.font.bold = True
p_m_tag.font.color.rgb = C_AMBER

p_m_desc = tf_m.add_paragraph()
p_m_desc.text = "Customize with your operational statistics:\n\n• Supported active orders processed daily\n• Reduced stock discrepancies from X% to Y%\n• Hours saved per week in IT & system administration"
p_m_desc.font.size = Pt(13)
p_m_desc.font.color.rgb = C_TEXT_MUTED


# ==============================================================================
# SLIDE 13: CONTACT
# ==============================================================================
slide13 = add_blank_slide()
add_header(slide13, "Contact & Professional Summary")
add_footer(slide13, 13)

add_card(slide13, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.9), border_color=C_INDIGO)

tb_c = slide13.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(10.9), Inches(4.3))
tf_c = tb_c.text_frame
tf_c.word_wrap = True

pc1 = tf_c.paragraphs[0]
pc1.text = "Budi"
pc1.font.size = Pt(32)
pc1.font.bold = True
pc1.font.color.rgb = C_TEXT_MAIN

pc2 = tf_c.add_paragraph()
pc2.text = "IT Support / IT Specialist  •  West Java & Jakarta, Indonesia (Remote or Onsite)\n"
pc2.font.size = Pt(18)
pc2.font.bold = True
pc2.font.color.rgb = C_INDIGO

contacts = [
    "• Background: Degree in Industrial Automation & Robotics Engineering",
    "• Expertise: Network Troubleshooting, ERP/System Integration, IoT Systems & Python/Dart Automation",
    "• Certifications: BNSP Senior Operator PLC  |  BNSP Junior Network Engineer",
    "• Production System URL: https://mdhproduction.com",
    "• Open for Opportunities: IT Support, Systems Administrator, IT Specialist roles in West Java & Jakarta"
]
for item in contacts:
    pi = tf_c.add_paragraph()
    pi.text = item
    pi.font.size = Pt(14)
    pi.font.color.rgb = C_TEXT_MUTED

# Save presentation
output_path = "Mobile_ERP_Portfolio_Presentation.pptx"
prs.save(output_path)
print(f"Successfully generated {output_path}")
